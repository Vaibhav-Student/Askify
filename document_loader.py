"""
document_loader.py — Document Processing Pipeline

Supports both:
1. In-memory streams (Flask / React UI) via load_pdf, load_pptx, load_xlsx, load_multiple_pdfs
2. File-path based chunked indexing (Streamlit app) via load_and_chunk
"""

import os
import re
import io
import logging
from pathlib import Path

from PyPDF2 import PdfReader
from pptx import Presentation
import openpyxl
from docx import Document as DocxDocument
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

from config import CHUNK_SIZE, CHUNK_OVERLAP, ALLOWED_EXTENSIONS

logger = logging.getLogger(__name__)


# ── In-Memory Stream Parsers (used by app_api.py) ────────────────────────────

def load_pdf(uploaded_file):
    """Extract text from an in-memory PDF file stream."""
    pdf_reader = PdfReader(uploaded_file)
    text = ""
    for page in pdf_reader.pages:
        extracted = page.extract_text()
        if extracted:
            text += extracted + "\n"
    return text


def load_pptx(uploaded_file):
    """Extract text from an in-memory PowerPoint presentation stream."""
    prs = Presentation(uploaded_file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)


def load_xlsx(uploaded_file):
    """Extract text from an in-memory Excel workbook stream."""
    wb = openpyxl.load_workbook(uploaded_file, data_only=True)
    text = []
    for sheet in wb.sheetnames:
        text.append(f"--- Sheet: {sheet} ---")
        ws = wb[sheet]
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
            if row_text.strip():
                text.append(row_text)
    return "\n".join(text)


def load_multiple_pdfs(uploaded_files):
    """Extract text from multiple in-memory PDF file streams."""
    all_text = ""
    for file in uploaded_files:
        text = load_pdf(file)
        all_text += text + "\n"
    return all_text


def extract_text_and_pages(file_bytes_or_stream, filename: str):
    """
    Extracts flat text and page-by-page structures from a file stream/bytes.
    Returns (text, pages) where pages is a list of {"text": str, "page": int/str}.
    """
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    text = ""
    pages = []
    
    # Ensure we have a stream/file-like object
    if isinstance(file_bytes_or_stream, bytes):
        stream = io.BytesIO(file_bytes_or_stream)
    else:
        stream = file_bytes_or_stream
        
    if ext == "pdf":
        try:
            pdf_reader = PdfReader(stream)
            for i, page in enumerate(pdf_reader.pages):
                extracted = page.extract_text()
                if extracted and extracted.strip():
                    extracted_strip = extracted.strip()
                    pages.append({"text": extracted_strip, "page": i + 1})
                    text += extracted_strip + "\n"
        except Exception as e:
            logger.error(f"Error parsing PDF: {e}")
            raise e
    elif ext in ("pptx", "ppt"):
        try:
            prs = Presentation(stream)
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    # Table shapes: extract cell text row by row
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_cells = []
                            for cell in row.cells:
                                cell_text = cell.text.strip()
                                if cell_text:
                                    row_cells.append(cell_text)
                            if row_cells:
                                slide_text.append(" | ".join(row_cells))
                    # Grouped shapes: recurse into group
                    elif shape.shape_type is not None and hasattr(shape, "shapes"):
                        for sub_shape in shape.shapes:
                            if hasattr(sub_shape, "text") and sub_shape.text.strip():
                                slide_text.append(sub_shape.text.strip())
                    # Regular shapes with text
                    elif hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                slide_content = "\n".join(slide_text).strip()
                if slide_content:
                    pages.append({"text": slide_content, "page": i + 1})
                    text += slide_content + "\n"
        except Exception as e:
            logger.error(f"Error parsing PPTX: {e}")
            raise e
    elif ext in ("xlsx", "xls"):
        try:
            wb = openpyxl.load_workbook(stream, data_only=True)
            for sheet in wb.sheetnames:
                sheet_text = []
                ws = wb[sheet]
                for row in ws.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        sheet_text.append(row_text.strip())
                sheet_content = "\n".join(sheet_text).strip()
                if sheet_content:
                    pages.append({"text": sheet_content, "page": sheet})
                    text += f"--- Sheet: {sheet} ---\n" + sheet_content + "\n"
        except Exception as e:
            logger.error(f"Error parsing Excel: {e}")
            raise e
    else:
        # Fallback for all other file types
        try:
            raw = stream.read() if hasattr(stream, "read") else stream
            if isinstance(raw, bytes):
                try:
                    text = raw.decode("utf-8")
                except Exception:
                    text = raw.decode("latin-1")
            else:
                text = str(raw)
        except Exception as e:
            logger.error(f"Error decoding text: {e}")
            text = f"[Binary/Non-text content for {filename}]"
            
        text_strip = text.strip()
        pages.append({"text": text_strip, "page": 1})
        text = text_strip
        
    return text, pages


# ── Document Type Detection & Path-Based Processing (used by app.py) ──────────

# Keyword patterns for classification
_DOC_TYPE_PATTERNS: dict[str, list[str]] = {
    "question_paper": [
        r"question\s*paper",
        r"examination",
        r"marks?\s*:",
        r"attempt\s+(any|all)",
        r"max\.?\s*marks",
        r"time\s*:\s*\d",
        r"answer\s+(any|all)",
        r"instructions\s+to\s+candidates",
    ],
    "lab_manual": [
        r"experiment\s*no",
        r"lab\s*(manual|report)",
        r"apparatus",
        r"procedure\s*:",
        r"observation\s*table",
        r"aim\s*:",
        r"theory\s*:",
        r"result\s*:",
    ],
    "notes": [
        r"chapter\s*\d",
        r"unit\s*\d",
        r"module\s*\d",
        r"lecture\s*notes",
        r"summary\s*:",
        r"key\s*points",
        r"introduction\s*:",
        r"definition\s*:",
    ],
}


def detect_document_type(text: str) -> str:
    """Classify a document based on keyword pattern matching."""
    sample = text[:2000].lower()
    scores: dict[str, int] = {}

    for doc_type, patterns in _DOC_TYPE_PATTERNS.items():
        score = sum(
            1 for pattern in patterns if re.search(pattern, sample, re.IGNORECASE)
        )
        scores[doc_type] = score

    best_type = max(scores, key=scores.get)
    return best_type if scores[best_type] >= 2 else "general"


def _extract_pdf(file_path: str) -> list[dict[str, str | int]]:
    """Extract text page-by-page from a PDF file."""
    reader = PdfReader(file_path)
    pages: list[dict[str, str | int]] = []

    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text and text.strip():
            pages.append({"text": text.strip(), "page": i + 1})

    if not pages:
        logger.warning("No text extracted from PDF: %s", file_path)

    return pages


def _extract_docx(file_path: str) -> list[dict[str, str | int]]:
    """Extract text from a DOCX file."""
    doc = DocxDocument(file_path)
    full_text = "\n".join(
        para.text.strip() for para in doc.paragraphs if para.text.strip()
    )

    if not full_text:
        logger.warning("No text extracted from DOCX: %s", file_path)
        return []

    return [{"text": full_text, "page": 1}]


def load_and_chunk(file_path: str) -> list[Document]:
    """Complete ingestion pipeline: extract → detect type → chunk → add metadata."""
    path = Path(file_path)

    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = path.suffix.lower()
    # Support docx and pdf, but allow extension matching
    if ext not in [".pdf", ".docx", ".doc"]:
        raise ValueError(
            f"Unsupported file type: {ext}. Allowed: pdf, docx, doc"
        )

    # Step 1: Extract raw text with page info
    if ext == ".pdf":
        pages = _extract_pdf(file_path)
    else:
        pages = _extract_docx(file_path)

    if not pages:
        return []

    # Step 2: Detect document type from first page's text
    full_text = " ".join(p["text"] for p in pages)
    doc_type = detect_document_type(full_text)
    logger.info("Detected document type for '%s': %s", path.name, doc_type)

    # Step 3: Create LangChain Documents with metadata
    raw_documents: list[Document] = []
    for page_data in pages:
        raw_documents.append(
            Document(
                page_content=str(page_data["text"]),
                metadata={
                    "source": path.name,
                    "doc_type": doc_type,
                    "page": page_data["page"],
                },
            )
        )

    # Step 4: Split into optimized chunks
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
        length_function=len,
        separators=["\n\n", "\n", ". ", " ", ""],
    )

    chunks = splitter.split_documents(raw_documents)

    # Step 5: Add chunk index to metadata for traceability
    for i, chunk in enumerate(chunks):
        chunk.metadata["chunk_index"] = i

    logger.info(
        "Processed '%s': %d pages → %d chunks (type: %s)",
        path.name,
        len(pages),
        len(chunks),
        doc_type,
    )

    return chunks
